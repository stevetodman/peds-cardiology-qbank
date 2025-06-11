import os
import subprocess
import tempfile
from io import BytesIO

DEFAULT_MERMAID = """graph TD
    A[Start] --> B{Decision}
    B -->|Yes| C[Do something]
    B -->|No| D[Do something else]
""".strip()


def demo_parser(guideline_text: str) -> str:
    """Return demo Mermaid code regardless of input."""
    return DEFAULT_MERMAID


def parse_guideline(guideline_text: str) -> str:
    """Generate mermaid code using OpenAI if available, else use demo parser."""
    api_key = os.getenv("OPENAI_API_KEY")
    if api_key:
        try:
            import openai  # type: ignore
        except Exception:
            return demo_parser(guideline_text)
        openai.api_key = api_key
        try:
            response = openai.ChatCompletion.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "Convert the following clinical guideline to a mermaid"
                            " flowchart. Return only the mermaid code."
                        ),
                    },
                    {"role": "user", "content": guideline_text},
                ],
                temperature=0,
            )
        except Exception:
            return demo_parser(guideline_text)
        return response.choices[0].message.content.strip()
    return demo_parser(guideline_text)


def export_png(mermaid_code: str) -> bytes:
    """Export Mermaid code to PNG using mermaid-cli."""
    with tempfile.TemporaryDirectory() as tmpdir:
        in_path = os.path.join(tmpdir, "diagram.mmd")
        out_path = os.path.join(tmpdir, "diagram.png")
        with open(in_path, "w") as f:
            f.write(mermaid_code)
        subprocess.run(["mmdc", "-i", in_path, "-o", out_path], check=True)
        with open(out_path, "rb") as f:
            return f.read()


def export_pptx(mermaid_code: str) -> bytes:
    """Export Mermaid code inside a PPTX slide."""
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches  # type: ignore

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    textbox.text = mermaid_code
    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()


def main() -> None:
    import streamlit as st  # imported here to avoid dependency during testing

    st.title("Guideline to Flowchart")
    guideline = st.text_area("Paste clinical guideline text", height=200)
    if st.button("Generate Flowchart") and guideline.strip():
        mermaid = parse_guideline(guideline)
        st.markdown(f"```mermaid\n{mermaid}\n```")
        if st.button("Create PNG"):
            try:
                png = export_png(mermaid)
                st.download_button("Download PNG", png, "flowchart.png")
            except Exception as exc:
                st.error(f"PNG export failed: {exc}")
        if st.button("Create PPTX"):
            pptx_data = export_pptx(mermaid)
            st.download_button("Download PPTX", pptx_data, "flowchart.pptx")


if __name__ == "__main__":
    main()
